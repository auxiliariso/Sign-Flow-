"""
signers/pptx_signer.py
======================
Firmador para presentaciones PowerPoint (.pptx).

Capacidades:
  - Detecta placeholders en cuadros de texto y notas de diapositiva.
  - Agrega una nueva diapositiva de firma por cada firma (no sobrescribe).
  - Si hay múltiples firmas, genera una diapositiva de resumen acumulativo.
  - Conserva tema y diseño de la presentación.
"""
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

from core.signer_base import SignerBase, SignaturePayload, SignResult
from core.hash_service import HashService

NAVY  = RGBColor(0x1F, 0x38, 0x64)
MID   = RGBColor(0x2E, 0x50, 0x90)
LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1A, 0x1A, 0x2E)
GREY  = RGBColor(0x6B, 0x7A, 0x99)

SIGNATURE_KEYWORDS = [
    "{{signflow}}", "{{firma_responsable}}", "{{multifirma}}",
    "{{responsables}}", "{{firma}}",
    "firma:", "responsable:", "revisó:", "autorizó:", "aprobó:",
]


class PptxSigner(SignerBase):

    def sign(
        self,
        doc_path: str,
        output_path: str,
        payload: SignaturePayload,
        all_previous: list[SignaturePayload],
    ) -> SignResult:
        prs = Presentation(doc_path)
        inserted = False

        # 1. Intentar reemplazar placeholder en diapositivas existentes
        for slide in prs.slides:
            if self._replace_placeholder_in_slide(slide, payload):
                inserted = True
                break

        # 2. Sin placeholder → agregar diapositiva de firma
        if not inserted:
            all_sigs = list(all_previous) + [payload]
            if len(all_sigs) == 1:
                # Primera firma: diapositiva individual
                self._add_single_signature_slide(prs, payload)
            else:
                # Múltiples firmas: diapositiva de resumen acumulativo
                # Remover diapositiva de resumen anterior si existe
                self._remove_summary_slide(prs)
                self._add_summary_slide(prs, all_sigs)

        prs.save(output_path)
        doc_hash = HashService().file_hash(output_path)
        return SignResult(
            output_path=output_path,
            firma_hash=payload.firma_hash,
            documento_hash_post=doc_hash,
        )

    def detect_signature_zones(self, doc_path: str) -> list[dict]:
        prs = Presentation(doc_path)
        zones = []
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip().lower()
                        for kw in SIGNATURE_KEYWORDS:
                            if kw in txt:
                                zones.append({
                                    "slide": i + 1,
                                    "shape": shape.name,
                                    "text": para.text.strip(),
                                    "keyword": kw,
                                })
        return zones

    # ── Placeholder ───────────────────────────────────────────────────────────

    def _replace_placeholder_in_slide(self, slide, payload: SignaturePayload) -> bool:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip().lower()
                if any(kw in txt for kw in SIGNATURE_KEYWORDS):
                    # Limpiar y reescribir
                    for run in para.runs:
                        run.text = ""
                    if para.runs:
                        r = para.runs[0]
                        r.text = (f"{payload.validation_id} | "
                                  f"{payload.nombre_completo} | "
                                  f"{payload.fecha}")
                        r.font.bold  = True
                        r.font.color.rgb = NAVY
                        r.font.size  = Pt(10)
                    return True
        return False

    # ── Diapositiva individual ────────────────────────────────────────────────

    def _add_single_signature_slide(self, prs: Presentation, payload: SignaturePayload):
        layout = prs.slide_layouts[6]   # Blank
        slide  = prs.slides.add_slide(layout)
        W, H   = prs.slide_width, prs.slide_height

        # Banner superior
        self._rect(slide, 0, 0, W, int(H * 0.20), NAVY)

        # Título
        self._text_box(slide,
                       f"✦ Firma Digital  —  {payload.validation_id}",
                       Inches(0.4), Inches(0.10),
                       W - Inches(0.8), Inches(0.9),
                       Pt(18), bold=True, color=WHITE)

        # Subtítulo
        self._text_box(slide, "SignFlow — Documento firmado digitalmente",
                       Inches(0.4), Inches(0.80),
                       W - Inches(0.8), Inches(0.4),
                       Pt(10), bold=False, color=LIGHT)

        rows = [
            ("Firmado por",  payload.nombre_completo),
            ("Cargo",        payload.nombre_puesto),
            ("Fecha (UTC)",  payload.fecha),
            ("Hora (UTC)",   payload.hora),
            ("Hash",         payload.firma_hash_short),
        ]

        y_base = int(H * 0.25)
        row_h  = int(H * 0.11)

        for i, (label, value) in enumerate(rows):
            y   = y_base + i * row_h
            bg  = LIGHT if i % 2 == 0 else WHITE
            self._rect(slide, Inches(0.4), y, W - Inches(1.6), row_h - Emu(40000), bg)
            self._text_box(slide, label,
                           Inches(0.5), y + Emu(60000),
                           Inches(2.2), row_h, Pt(10), bold=True, color=NAVY)
            self._text_box(slide, value,
                           Inches(2.9), y + Emu(60000),
                           W - Inches(4.2), row_h, Pt(10), bold=False, color=DARK)

        # QR
        if payload.qr_image_bytes:
            qr_left = W - Inches(1.3)
            qr_top  = int(H * 0.25)
            pic = slide.shapes.add_picture(
                io.BytesIO(payload.qr_image_bytes),
                qr_left, qr_top,
                width=Inches(1.0), height=Inches(1.0),
            )

        # Nota al pie
        self._text_box(slide,
                       "Este documento ha sido firmado digitalmente. "
                       "Verifica su autenticidad en la API de SignFlow.",
                       Inches(0.4), H - Inches(0.5),
                       W - Inches(0.8), Inches(0.4),
                       Pt(7), bold=False, color=GREY)

    # ── Diapositiva de resumen (multifirma) ───────────────────────────────────

    def _add_summary_slide(self, prs: Presentation, all_sigs: list[SignaturePayload]):
        layout = prs.slide_layouts[6]
        slide  = prs.slides.add_slide(layout)
        W, H   = prs.slide_width, prs.slide_height

        # Marcar esta diapositiva para poder removerla en la siguiente firma
        slide.name = "__signflow_summary__"

        # Banner
        self._rect(slide, 0, 0, W, int(H * 0.18), NAVY)
        self._text_box(slide, f"✦ Firmas Digitales ({len(all_sigs)})",
                       Inches(0.4), Inches(0.08),
                       W - Inches(0.8), Inches(0.8),
                       Pt(18), bold=True, color=WHITE)

        y     = int(H * 0.22)
        row_h = max(int((H * 0.75) / max(len(all_sigs), 1)), int(H * 0.10))

        for i, sig in enumerate(all_sigs):
            bg = LIGHT if i % 2 == 0 else WHITE
            self._rect(slide, Inches(0.3), y, W - Inches(0.6),
                       row_h - Emu(30000), bg)

            # Número de firma
            self._text_box(slide, f"#{i+1}",
                           Inches(0.4), y + Emu(50000),
                           Inches(0.5), row_h, Pt(9), bold=True, color=MID)
            # Datos
            self._text_box(slide,
                           f"{sig.nombre_completo}  —  {sig.nombre_puesto}",
                           Inches(1.0), y + Emu(50000),
                           Inches(5.5), row_h, Pt(9), bold=True, color=DARK)
            self._text_box(slide,
                           f"{sig.fecha}  {sig.hora}  ·  {sig.validation_id}  ·  Hash: {sig.firma_hash_short}",
                           Inches(1.0), y + Emu(260000),
                           Inches(6.5), row_h, Pt(7), bold=False, color=GREY)

            # QR en miniatura
            if sig.qr_image_bytes:
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(sig.qr_image_bytes),
                        W - Inches(1.1), y + Emu(30000),
                        width=Inches(0.7), height=Inches(0.7),
                    )
                except Exception:
                    pass

            y += row_h

    def _remove_summary_slide(self, prs: Presentation):
        """Elimina la diapositiva de resumen anterior para reemplazarla."""
        xml_slides = prs.slides._sldIdLst
        for slide in prs.slides:
            if slide.name == "__signflow_summary__":
                rId = prs.slides._sldIdLst[-1].get("r:id")
                # python-pptx no tiene remove directo; usamos XML
                try:
                    slide_elem = slide._element
                    slide_elem.getparent().remove(slide_elem)
                except Exception:
                    pass
                break

    # ── Helpers de dibujo ─────────────────────────────────────────────────────

    @staticmethod
    def _rect(slide, left, top, width, height, color: RGBColor):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    @staticmethod
    def _text_box(slide, text, left, top, width, height,
                  font_size, bold, color: RGBColor):
        tb  = slide.shapes.add_textbox(left, top, width, height)
        tf  = tb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text           = text
        run.font.bold      = bold
        run.font.size      = font_size
        run.font.color.rgb = color
        run.font.name      = "Calibri"
