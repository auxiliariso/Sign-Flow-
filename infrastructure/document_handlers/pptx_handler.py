"""
Handler - Firma en presentaciones PowerPoint (.pptx)
Agrega una diapositiva de firma al final de la presentación.
"""
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from core.interfaces.repositories import IDocumentSigner

BRAND = RGBColor(0x1F, 0x38, 0x64)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xD6, 0xE4, 0xF0)


class PptxSigner(IDocumentSigner):
    def sign(
        self,
        doc_path: str,
        output_path: str,
        nombre_completo: str,
        nombre_puesto: str,
        firma_hash: str,
        fecha: str,
        hora: str,
    ) -> str:
        shutil.copy2(doc_path, output_path)
        prs = Presentation(output_path)

        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        W = prs.slide_width
        H = prs.slide_height

        # Fondo del banner superior
        _add_rect(slide, 0, 0, W, int(H * 0.22), BRAND)

        # Título
        _add_text_box(
            slide,
            text="✦ Documento Firmado Digitalmente",
            left=Inches(0.4), top=Inches(0.15),
            width=W - Inches(0.8), height=Inches(0.8),
            font_size=Pt(22), bold=True, color=WHITE,
        )

        data = [
            ("Firmado por", nombre_completo),
            ("Cargo / Puesto", nombre_puesto),
            ("Fecha (UTC)", fecha),
            ("Hora (UTC)", hora),
            ("Hash de verificación", firma_hash),
        ]
        top_base = int(H * 0.30)
        row_h = int(H * 0.12)

        for idx, (label, value) in enumerate(data):
            y = top_base + idx * row_h
            bg = LIGHT if idx % 2 == 0 else WHITE
            _add_rect(slide, Inches(0.3), y, W - Inches(0.6), row_h - Emu(30000), bg)
            _add_text_box(slide, label, Inches(0.4), y + Emu(50000),
                          Inches(2.5), row_h, Pt(11), True, BRAND)
            _add_text_box(slide, value, Inches(3.1), y + Emu(50000),
                          W - Inches(3.5), row_h, Pt(10), False, RGBColor(0x33, 0x33, 0x33))

        prs.save(output_path)
        return output_path


def _add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()


def _add_text_box(slide, text, left, top, width, height, font_size, bold, color):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
